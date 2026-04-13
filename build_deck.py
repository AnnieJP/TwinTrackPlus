from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # off-white body bg
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # dark text
MID    = RGBColor(0x44, 0x55, 0x6E)   # secondary text
ACCENT2= RGBColor(0x90, 0xE0, 0xEF)   # soft teal highlight

def hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, l, t, w, h, size=14, color=DARK, title=None, title_color=None):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color or TEAL
        first = False
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

# ════════════════════════════════════════════════════════════════
#  PRESENTATION DECK
# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # completely blank

# ── Slide 1: Title ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 5.8, 13.33, 1.7, TEAL)
# accent bar
add_rect(s, 0.5, 2.6, 0.12, 1.8, TEAL)
add_text(s, "TwinTrack", 0.75, 2.5, 9, 1.4, size=60, bold=True, color=WHITE)
add_text(s, "Digital Twin Economic Simulator for Small Business", 0.75, 3.85, 10, 0.7,
         size=22, color=ACCENT2)
add_text(s, "WEHack UTD  ·  Capital One Track  ·  April 2026",
         0.75, 6.0, 10, 0.6, size=16, color=NAVY, italic=True)

# ── Slide 2: The Problem ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT)
add_rect(s, 0, 0, 13.33, 1.2, NAVY)
add_rect(s, 0, 1.2, 0.18, 6.3, TEAL)
add_text(s, "THE PROBLEM", 0.4, 0.18, 10, 0.7, size=13, bold=True, color=TEAL)
add_text(s, "Every decision is a risk.", 0.4, 0.6, 10, 0.55, size=28, bold=True, color=WHITE)

# big stat box
add_rect(s, 0.5, 1.6, 5.8, 2.4, NAVY)
add_text(s, "82%", 0.7, 1.65, 5.4, 1.2, size=64, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, "of small business failures are linked\nto poor financial decisions",
         0.7, 2.75, 5.4, 0.9, size=14, color=WHITE, align=PP_ALIGN.CENTER)

bullets = [
    'No sandbox to test: "What if I raise prices 15%?"',
    'No way to model: "Should I open a second location?"',
    'No data on: "Which audience should I target?"',
    "Enterprise tools (Anaplan, PharmaComp) are out of reach — 6-figure licenses, months of setup",
    "Small business owners are flying blind on decisions that could end their business",
]
add_bullet_box(s, bullets, 6.8, 1.5, 6.0, 5.5, size=14, color=DARK)

# ── Slide 3: The Solution ────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT)
add_rect(s, 0, 0, 13.33, 1.2, NAVY)
add_rect(s, 0, 1.2, 0.18, 6.3, TEAL)
add_text(s, "THE SOLUTION", 0.4, 0.18, 10, 0.7, size=13, bold=True, color=TEAL)
add_text(s, "A risk-free sandbox for your business.", 0.4, 0.6, 10, 0.55, size=26, bold=True, color=WHITE)

add_text(s, "Register your business as a Digital Twin —\na live structured snapshot of your financials,\ncosts, and sales profile.",
         0.5, 1.4, 5.5, 2.0, size=17, color=DARK)

add_text(s, "Then simulate any decision before it costs you.", 0.5, 3.35, 5.5, 0.8, size=17,
         bold=True, color=NAVY)

use_cases = [
    ("Pricing Changes", "Raise or lower average price — demand modelled via CPI-derived elasticity"),
    ("Audience Shift",  "Target a new demographic — ticket size, footfall, and marketing costs modelled"),
    ("Franchise Expansion", "Open new locations — upfront costs amortized over 36 months with break-even projection"),
]
for i, (title, desc) in enumerate(use_cases):
    top = 1.5 + i * 1.65
    add_rect(s, 6.5, top, 6.3, 1.5, NAVY)
    add_rect(s, 6.5, top, 0.12, 1.5, TEAL)
    add_text(s, title, 6.75, top + 0.1, 5.8, 0.5, size=15, bold=True, color=TEAL)
    add_text(s, desc,  6.75, top + 0.55, 5.8, 0.85, size=13, color=WHITE)

# ── Slide 4: How It Works ────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT)
add_rect(s, 0, 0, 13.33, 1.2, NAVY)
add_text(s, "HOW IT WORKS", 0.4, 0.18, 10, 0.7, size=13, bold=True, color=TEAL)
add_text(s, "Four-stage pipeline. One decision.", 0.4, 0.6, 12, 0.55, size=26, bold=True, color=WHITE)

stages = [
    ("01", "Live Data Ingestion",
     "FRED · BLS · BEA · Census · NewsData.io\nDate-keyed cache prevents redundant API calls"),
    ("02", "Market Snapshot (MS)",
     "Claude Haiku resolves NAICS + MSA codes\nComputes CPI elasticity, sentiment, ARIMA 12-mo forecasts"),
    ("03", "Simulation Engine",
     "IP1 (control) vs IP2 (experiment)\nFixed/variable cost split · franchise amortization · audience modelling"),
    ("04", "Results Dashboard",
     "OP1 vs OP2: revenue, margin, footfall delta\nConfidence score · sentiment ribbon · plain-English verdict"),
]

for i, (num, title, desc) in enumerate(stages):
    col = i * 3.2 + 0.35
    add_rect(s, col, 1.5, 2.9, 5.5, NAVY)
    add_rect(s, col, 1.5, 2.9, 0.12, TEAL)
    add_text(s, num,   col+0.15, 1.65, 2.6, 0.9, size=38, bold=True, color=TEAL)
    add_text(s, title, col+0.15, 2.5,  2.6, 0.6, size=14, bold=True, color=WHITE)
    add_text(s, desc,  col+0.15, 3.15, 2.6, 3.5, size=12, color=ACCENT2)

# arrows between boxes
for i in range(3):
    ax = 0.35 + i*3.2 + 2.9 + 0.05
    add_text(s, "→", ax, 3.5, 0.25, 0.5, size=20, bold=True, color=TEAL)

# ── Slide 5: Tech Stack ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT)
add_rect(s, 0, 0, 13.33, 1.2, NAVY)
add_text(s, "TECH STACK", 0.4, 0.18, 10, 0.7, size=13, bold=True, color=TEAL)
add_text(s, "Lightweight. Fast. Production-ready.", 0.4, 0.6, 12, 0.55, size=26, bold=True, color=WHITE)

stack = [
    ("Frontend",      "React 19 + Vite 5",               "Component-driven SPA — no router, minimal overhead, sub-600ms load"),
    ("Backend",       "Python ThreadingHTTPServer",        "Zero-framework REST API — concurrent requests, 6 endpoints, pure stdlib"),
    ("Data Layer",    "5 Federal APIs + disk cache",       "FRED, BLS, BEA, Census, NewsData.io — date-keyed cache for reproducibility"),
    ("AI/LLM",        "Claude Haiku (claude-haiku-4-5)",   "NAICS resolution · NL→parameter extraction · grounded recommendation generation"),
    ("Forecasting",   "statsmodels + pmdarima ARIMA",      "Auto-order selection · 12-month revenue & margin projections per simulation"),
]

for i, (layer, tech, desc) in enumerate(stack):
    top = 1.45 + i * 1.12
    add_rect(s, 0.4, top, 12.5, 1.0, NAVY)
    add_rect(s, 0.4, top, 0.12, 1.0, TEAL)
    add_text(s, layer, 0.65, top+0.05, 2.2, 0.45, size=11, bold=True, color=TEAL)
    add_text(s, tech,  0.65, top+0.47, 2.2, 0.45, size=13, bold=True, color=WHITE)
    add_text(s, desc,  3.1,  top+0.15, 9.6, 0.7,  size=13, color=ACCENT2)

# ── Slide 6: What We're Proud Of ────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT)
add_rect(s, 0, 0, 13.33, 1.2, NAVY)
add_text(s, "WHAT WE'RE PROUD OF", 0.4, 0.18, 12, 0.7, size=13, bold=True, color=TEAL)
add_text(s, "Addressing a real gap — not a toy demo.", 0.4, 0.6, 12, 0.55, size=26, bold=True, color=WHITE)

add_text(s, "We looked at the state of the art:",
         0.5, 1.4, 12, 0.5, size=15, bold=True, color=DARK)

competitors = ["Anaplan", "Marketplace Simulation", "PharmaComp"]
for i, c in enumerate(competitors):
    add_rect(s, 0.5 + i*3.3, 2.0, 3.0, 0.7, NAVY)
    add_text(s, c, 0.5 + i*3.3, 2.05, 3.0, 0.6, size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

differentiators = [
    "Live macro-economic data — not static assumptions. CPI, unemployment, GDP updated in real time.",
    "Hyper-local market context — NAICS + MSA resolution means a Dallas bakery gets Dallas data, not national averages.",
    "Natural language input — describe your decision in plain English; the LLM extracts the simulation parameters.",
    "Explainable output — every result includes the underlying delta numbers, a confidence score, and a 2-sentence verdict.",
    "Graceful degradation — all LLM calls fail silently; the simulation still runs on pure economic data.",
]
add_bullet_box(s, differentiators, 0.5, 2.9, 12.3, 4.2, size=14, color=DARK)

# ── Slide 7: Challenges ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT)
add_rect(s, 0, 0, 13.33, 1.2, NAVY)
add_text(s, "CHALLENGES", 0.4, 0.18, 12, 0.7, size=13, bold=True, color=TEAL)
add_text(s, "Turning ambition into a working system.", 0.4, 0.6, 12, 0.55, size=26, bold=True, color=WHITE)

challenges = [
    ("Scope vs. POC",
     "We designed a full product — multi-tenant, persistent digital twins, scenario branching. Scoping it to a compelling hackathon POC without gutting the architecture was the hardest call we made."),
    ("Cross-expertise integration",
     "Four people, four layers: frontend, ML pipeline, simulation engine, and LLM layer. Each component worked in isolation. Making them speak the same data contract — IP1/IP2 JSON — required constant re-alignment."),
    ("Live API reliability",
     "Five federal APIs, each with different rate limits, response formats, and occasional outages. We built a date-keyed disk cache so a single bad API call never kills a full simulation run."),
    ("LLM + deterministic sim",
     "Claude Haiku drives context resolution and recommendations, but the financial simulation must be deterministic. Keeping the boundary clean — LLM enriches, engine computes — was a deliberate architectural decision."),
]

for i, (title, desc) in enumerate(challenges):
    top = 1.45 + i * 1.45
    add_rect(s, 0.4, top, 12.5, 1.3, NAVY)
    add_rect(s, 0.4, top, 0.12, 1.3, TEAL)
    add_text(s, title, 0.65, top+0.05, 3.5, 0.45, size=13, bold=True, color=TEAL)
    add_text(s, desc,  0.65, top+0.5,  11.8, 0.75, size=12, color=WHITE)

# ── Slide 8: Future Plans ────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 6.5, 13.33, 1.0, TEAL)
add_text(s, "WHAT'S NEXT", 0.5, 0.3, 12, 0.6, size=13, bold=True, color=TEAL)
add_text(s, "TwinTrack v2", 0.5, 0.85, 12, 1.1, size=52, bold=True, color=WHITE)

roadmap = [
    ("Educational Module",
     "Onboarding wizard for new businesses — guided setup with industry benchmarks and starter financial templates."),
    ("Richer Financial Models",
     "Move beyond common market assumptions. Incorporate seasonal demand curves, regional wage data, and sector-specific cost structures."),
    ("API Scaling",
     "Expand data ingestion to World Bank, local property indexes, and social media sentiment. Increase simulation resolution from monthly to weekly."),
    ("Multi-scenario Branching",
     "Run multiple experiments in parallel — compare 3 pricing strategies simultaneously, not just control vs. one experiment."),
]

for i, (title, desc) in enumerate(roadmap):
    col = (i % 2) * 6.4 + 0.5
    top = 2.3 + (i // 2) * 2.2
    add_rect(s, col, top, 6.0, 1.9, RGBColor(0x16, 0x2A, 0x42))
    add_rect(s, col, top, 0.12, 1.9, TEAL)
    add_text(s, title, col+0.25, top+0.1,  5.5, 0.5, size=14, bold=True, color=TEAL)
    add_text(s, desc,  col+0.25, top+0.6,  5.5, 1.2, size=12, color=ACCENT2)

add_text(s, "Built at WEHack UTD · Capital One Track · April 2026",
         0.5, 6.6, 12, 0.5, size=13, color=NAVY, align=PP_ALIGN.CENTER)

deck_path = "C:/Users/neeha/Desktop/WeHack vFinal/TwinTrack_Deck.pptx"
prs.save(deck_path)
print(f"Deck saved: {deck_path}")


# ════════════════════════════════════════════════════════════════
#  POSTER  (landscape 36×24)
# ════════════════════════════════════════════════════════════════
poster = Presentation()
W, H = 36, 24
poster.slide_width  = Inches(W)
poster.slide_height = Inches(H)
blank = poster.slide_layouts[6]

s = poster.slides.add_slide(blank)

# Background
add_rect(s, 0, 0, W, H, NAVY)

# ── Header ────────────────────────────────────────────────────────
add_rect(s, 0, 0, W, 3.8, RGBColor(0x06, 0x0F, 0x1A))
add_rect(s, 0, 3.8, W, 0.18, TEAL)
add_text(s, "TwinTrack", 0.8, 0.2, 18, 2.0, size=80, bold=True, color=WHITE)
add_text(s, "Digital Twin Economic Simulator for Small Business",
         0.8, 2.2, 22, 1.0, size=26, color=ACCENT2)
add_text(s, "WEHack UTD  ·  Capital One Track  ·  April 2026",
         0.8, 3.1, 22, 0.6, size=18, color=MID, italic=True)

# ── 3-column body ─────────────────────────────────────────────────
TOP = 4.2
BOT = 22.8
CH  = BOT - TOP   # 18.6 inches tall
CW  = 11.0

# Column positions
C1, C2, C3 = 0.4, 12.5, 24.6

# ── Column 1: Problem + Solution ─────────────────────────────────
# Problem
PH = 8.5
add_rect(s, C1, TOP, CW, PH, RGBColor(0x10, 0x22, 0x35))
add_rect(s, C1, TOP, 0.18, PH, TEAL)
add_text(s, "THE PROBLEM", C1+0.35, TOP+0.2, CW-0.5, 0.7, size=15, bold=True, color=TEAL)
add_text(s, "Small businesses have no\nrisk-free environment\nto test decisions.",
         C1+0.35, TOP+1.0, CW-0.5, 2.5, size=22, bold=True, color=WHITE)
p_items = [
    "Pricing, expansion, audience shifts — all carry real financial risk",
    "Enterprise simulators cost 6-figures, months to deploy",
    "Owners rely on gut instinct instead of data",
    "One bad decision can end a small business",
]
add_bullet_box(s, p_items, C1+0.35, TOP+3.7, CW-0.5, 4.5, size=16, color=ACCENT2)

# Solution
ST = TOP + PH + 0.4
SH = CH - PH - 0.4
add_rect(s, C1, ST, CW, SH, RGBColor(0x10, 0x22, 0x35))
add_rect(s, C1, ST, 0.18, SH, TEAL)
add_text(s, "THE SOLUTION", C1+0.35, ST+0.2, CW-0.5, 0.7, size=15, bold=True, color=TEAL)
add_text(s, "Register. Simulate. Decide.",
         C1+0.35, ST+1.0, CW-0.5, 1.0, size=22, bold=True, color=WHITE)
s_items = [
    "Register your business as a Digital Twin",
    "Pick a use case: Pricing · Audience · Franchise",
    "Describe your decision in plain English",
    "Run the engine — OP1 (control) vs OP2 (experiment)",
    "Read the verdict: revenue delta, margin, confidence score",
]
add_bullet_box(s, s_items, C1+0.35, ST+2.1, CW-0.5, 6.5, size=16, color=ACCENT2)

# ── Column 2: Pipeline ────────────────────────────────────────────
add_rect(s, C2, TOP, CW, CH, RGBColor(0x0A, 0x18, 0x28))
add_rect(s, C2, TOP, 0.18, CH, TEAL)
add_text(s, "PIPELINE", C2+0.35, TOP+0.2, CW-0.5, 0.7, size=15, bold=True, color=TEAL)
add_text(s, "Four stages. One decision.",
         C2+0.35, TOP+1.0, CW-0.5, 0.8, size=20, bold=True, color=WHITE)

stages_p = [
    ("01  Data Ingestion",
     "FRED · BLS · BEA · Census · NewsData.io\nDate-keyed disk cache prevents redundant calls"),
    ("02  Market Snapshot",
     "Claude Haiku resolves NAICS + MSA codes\nCPI elasticity · ARIMA 12-month forecasts · sentiment"),
    ("03  Simulation Engine",
     "IP1 (control) vs IP2 (experiment)\nFixed/variable cost split · franchise amortization"),
    ("04  Results Dashboard",
     "Revenue & margin delta · confidence score\nSentiment ribbon · plain-English verdict"),
]
for i, (title, desc) in enumerate(stages_p):
    bt = TOP + 2.1 + i * 4.1
    add_rect(s, C2+0.35, bt, CW-0.5, 3.7, RGBColor(0x0D, 0x1B, 0x2A))
    add_rect(s, C2+0.35, bt, CW-0.5, 0.15, TEAL)
    add_text(s, title, C2+0.55, bt+0.25, CW-0.8, 0.7, size=17, bold=True, color=TEAL)
    add_text(s, desc,  C2+0.55, bt+1.05, CW-0.8, 2.5, size=15, color=WHITE)
    if i < 3:
        add_text(s, "▼", C2+5.0, bt+3.7, 1.0, 0.5, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ── Column 3: Tech Stack + Why TwinTrack ─────────────────────────
TSH = 9.5
add_rect(s, C3, TOP, CW, TSH, RGBColor(0x10, 0x22, 0x35))
add_rect(s, C3, TOP, 0.18, TSH, TEAL)
add_text(s, "TECH STACK", C3+0.35, TOP+0.2, CW-0.5, 0.7, size=15, bold=True, color=TEAL)

ts = [
    ("Frontend",    "React 19 + Vite 5"),
    ("Backend",     "Python ThreadingHTTPServer"),
    ("Data",        "5 Federal APIs + disk cache"),
    ("AI / LLM",    "Claude Haiku (Anthropic)"),
    ("Forecasting", "statsmodels ARIMA + pmdarima"),
]
for i, (layer, tech) in enumerate(ts):
    rt = TOP + 1.2 + i * 1.6
    add_rect(s, C3+0.35, rt, CW-0.5, 1.4, RGBColor(0x0D, 0x1B, 0x2A))
    add_text(s, layer, C3+0.55, rt+0.1, 3.5, 0.5, size=13, bold=True, color=TEAL)
    add_text(s, tech,  C3+0.55, rt+0.65, CW-0.8, 0.6, size=15, color=WHITE)

# Why TwinTrack
WT = TOP + TSH + 0.4
WH = CH - TSH - 0.4
add_rect(s, C3, WT, CW, WH, RGBColor(0x10, 0x22, 0x35))
add_rect(s, C3, WT, 0.18, WH, TEAL)
add_text(s, "WHY TWINTRACK", C3+0.35, WT+0.2, CW-0.5, 0.7, size=15, bold=True, color=TEAL)
diffs = [
    "Live macro data — CPI, GDP, unemployment in real time",
    "Hyper-local — NAICS+MSA means city-level market context",
    "NL input — plain English description, LLM extracts params",
    "Explainable — delta numbers + confidence + verdict",
    "Graceful degradation — works without LLM if key is missing",
]
add_bullet_box(s, diffs, C3+0.35, WT+1.1, CW-0.5, WH-1.5, size=16, color=ACCENT2)

# ── Footer ────────────────────────────────────────────────────────
add_rect(s, 0, 23.0, W, 1.0, TEAL)
add_text(s, "TwinTrack  ·  WEHack UTD  ·  Capital One Track  ·  April 2026",
         0, 23.2, W, 0.6, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

poster_path = "C:/Users/neeha/Desktop/WeHack vFinal/TwinTrack_Poster.pptx"
poster.save(poster_path)
print(f"Poster saved: {poster_path}")
